import google.generativeai as genai
import numpy as np
from sentence_transformers import SentenceTransformer
import faiss
from dotenv import load_dotenv
import os
import re
import traceback

# for api
load_dotenv() 
api_key = os.getenv("GEMINI_API_KEY")

# check if API key is loaded
if not api_key:
    print("Error: GEMINI_API_KEY not found in .env file")
    print("Current working directory:", os.getcwd())
    print("Files in directory:", os.listdir("."))
    exit(1)

# configure the API key
genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-1.5-flash")

def pdf_ocr(path):
    from pdf2image import convert_from_path
    import pytesseract
    try:
        images = convert_from_path(path)

        full_text = f"File: {os.path.basename(path)} contents: after completing OCR: " + '-'*20 + '\n'

        for i, image in enumerate(images):
            text = pytesseract.image_to_string(image)
            full_text += text + "\n\n"
        
        print(f"OCR completed successfully.")
        return full_text
    
    except Exception as e:
        print(f"An error occurred during OCR: {e}")
        return ""

def read_pdf(path, use_ocr=False):
    if not os.path.exists(path):
        print(f"Error: PDF file '{path}' not found!")
        return ""
    
    text = f"File: {os.path.basename(path)} contents: " + '-'*20 + '\n'
    
    if not use_ocr:
        try:
            import pypdf
            with open(path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                print(f"PDF has {len(pdf_reader.pages)} pages")
                
                for i, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text() or ""
                    print(f"Page {i+1} extracted {len(page_text)} characters")
                    if page_text.strip():
                        text += page_text + "\n"
                    else:
                        print(f"Warning: Page {i+1} appears to be empty or image-only")
        
        except Exception as e:
            print(f"Error reading PDF with pypdf: {e}")
            text = ""
        
        print(f"pypdf extraction completed. Total text length: {len(text)} characters")
        
        if len(text.strip()) < 100:
            print("Very little text extracted with pypdf. This might be an image-based PDF.")
            print("Automatically switching to OCR...")
            return pdf_ocr(path)
        else:
            return text
    else:
        return pdf_ocr(path)

def read_docx(path):
    if not os.path.exists(path):
        print(f"Error: Word document file '{path}' not found!")
        return ""

    try:
        from docx import Document
        document = Document(path)
        text = "\n".join([para.text for para in document.paragraphs])
        text = f"File: {os.path.basename(path)} contents: \n" + "-"*20 + "\n" + text + "\n" + "-"*20
        print(f"docx extraction completed. Total text length: {len(text)} characters")
        return text
    except Exception as e:
        print(f"Error reading docx file: {e}")
        return ""

def read_excel(path, csv):
    if not os.path.exists(path):
        print(f"Error: Spreadsheet file '{path}' not found!")
        return ""

    try:
        import pandas as pd
        text = f"File: {os.path.basename(path)} contents: \n" + "-"*20 + "\n"

        if csv == False:
            excel_data = pd.read_excel(path, sheet_name=None)
            for sheet_name, df in excel_data.items():
                text += f"\n-----Sheet Name: {sheet_name}:-----\n"
                text += df.to_string(index=False)
                text += '\n' + '-'*20 + '\n'
            return text
        else:
            csv_data = pd.read_csv(path)
            text += csv_data.to_string(index=False)
            text += '\n' + '-'*20
            return text
    except Exception as e:
        print(f"Error reading spreadsheet file: {e}")
        return ""

def read_text(path):
    if not os.path.exists(path):
        print(f"Error: text file '{path}' not found!")
        return ""

    text = f"File: {os.path.basename(path)} contents: \n" + "-"*20 + "\n"

    try:
        with open(path, "r", encoding="utf-8") as f:
            text += f.read()
        print(f"text extraction completed. Total text length: {len(text)} characters")
        return text
    except Exception as e:
        print(f"Error reading text file: {e}")
        return ""

def read_pptx(path):
    if not os.path.exists(path):
        print(f"Error: PowerPoint file '{path}' not found!")
        return ""

    try:
        from pptx import Presentation
        text = f"File: {os.path.basename(path)} contents: \n" + "-"*20 + "\n"
        presentation = Presentation(path)
        for i, slide in enumerate(presentation.slides):
            text += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text.strip() + "\n"
        text += "-" * 20
        print(f"pptx extraction completed. Total text length: {len(text)} characters")
        return text
    except Exception as e:
        print(f"Error reading PowerPoint file: {e}")
        return ""

def split_into_chunks(text, max_chunk_size=1000):
    if not text or not text.strip():
        print("Warning: No text to split into chunks!")
        return []
    
    sentences = re.split(r'(?<=[.!?]) +', text)
    print(f"Split text into {len(sentences)} sentences")
    
    chunks, chunk = [], ""
    for sentence in sentences:
        if len(chunk) + len(sentence) < max_chunk_size:
            chunk += sentence + " "
        else:
            if chunk.strip(): 
                chunks.append(chunk.strip())
            chunk = sentence + " "
    
    if chunk and chunk.strip():
        chunks.append(chunk.strip())
    
    print(f"Created {len(chunks)} chunks")
    return chunks

def embed_chunks(chunks):
    embedder = SentenceTransformer('all-MiniLM-L6-v2')
    embeddings = embedder.encode(chunks)
    embeddings = np.array(embeddings).astype('float32')
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    return index, embedder

def get_top_k_chunks(query, embedder, index, chunks, k=3):
    query_vector = embedder.encode([query]).astype('float32')
    D, I = index.search(query_vector, k)
    return [chunks[i] for i in I[0]]

def ask_gemini(query, context):
    prompt = f"""Based on the following document content, answer the user's question.

---Document Content---
{context}

---Question---
{query}
"""
    return get_response(prompt)

def get_response(message):
    try:
        response = model.generate_content(message)
        return response.text if hasattr(response, 'text') else response.candidates[0].content.parts[0].text
    except Exception as e:
        print(f"Error getting response from Gemini: {e}")
        return "Sorry, I couldn't generate a response. Please try again."

def main():
    print("=" * 60)
    print("         DOCUMENT CHAT ASSISTANT")
    print("=" * 60)
    print("Commands:")
    print("- Type your question to ask about the document")
    print("- Type 'load' to load a new document file")
    print("- Type 'exit' or 'quit' to stop")
    print("=" * 60)

    chunks = []
    index = None
    embedder = None
    current_file = None

    def load_document():
        nonlocal chunks, index, embedder, current_file

        while True:
            file_path = input("\nEnter file path or 'cancel' to go back (supported formats: .pdf .docx .xlsx .csv .txt .pptx): ").strip()
            
            if file_path.lower() == 'cancel':
                return False
            
            if not os.path.exists(file_path):
                print(f"Error: File '{file_path}' not found.")
                continue
            
            print(f"\n{'='*50}\n")
            print(f"Loading File: {file_path}")
            print(f"\n{'='*50}")

            try:
                file_name = os.path.basename(file_path)
                ext = os.path.splitext(file_path)[1].lower()
                text = ""

                if ext == ".pdf":
                    text = read_pdf(file_path, use_ocr=False)
                elif ext == ".docx":
                    text = read_docx(file_path)
                elif ext == ".xlsx":
                    text = read_excel(file_path, csv=False)
                elif ext == ".csv":
                    text = read_excel(file_path, csv=True)
                elif ext == ".txt":
                    text = read_text(file_path)
                elif ext == ".pptx":
                    text = read_pptx(file_path)
                else:
                    print(f"Unsupported file format: {ext}")
                    continue
                
                if not text or not text.strip():
                    print("ERROR: No text extracted from file!")
                    continue
                
                chunks = split_into_chunks(text)
                if not chunks:
                    print("ERROR: No chunks created from text!")
                    continue

                index, embedder = embed_chunks(chunks)
                current_file = file_path

                print(f"\n✅ File loaded successfully!")
                print(f"📄 File: {file_name}")
                print(f"📊 Text length: {len(text):,} characters")
                print(f"📋 Chunks created: {len(chunks)}")
                print("\nYou can now ask questions about this document!")
                return True

            except Exception as e:
                print(f"Error loading file: {e}")
                traceback.print_exc()
                continue

    if not load_document():
        print("No document loaded. Exiting...")
        return

    print(f"\n{'='*60}")
    print(" CHAT MODE - Ask questions about your document!")
    print(f"{'='*60}")
    
    while True:
        try:
            user_input = input(f"\nYou (File: {os.path.basename(current_file) if current_file else 'None'}): ").strip()
            
            if not user_input:
                continue
            
            if user_input.lower() in ['exit', 'quit', 'q']:
                print("Goodbye!")
                break
            elif user_input.lower() == 'load':
                load_document()
                continue
            elif user_input.lower() in ['help', 'h']:
                print("\nAvailable commands:")
                print("- Ask any question about the loaded document")
                print("- 'load' - Load a new document file")
                print("- 'help' - Show this help message")
                print("- 'exit' or 'quit' - Exit the program")
                continue
            
            if not chunks or not index or not embedder:
                print("No document loaded! Use 'load' command to load a document first.")
                continue
            
            print("🔍 Searching relevant content...")
            relevant_chunks = get_top_k_chunks(user_input, embedder, index, chunks, k=5)
            context = "\n".join(relevant_chunks)
            
            if not context:
                print("No relevant content found for your question.")
                continue
            
            print("🤔 Thinking...")
            answer = ask_gemini(user_input, context)
            print(f"\nAssistant: {answer}")
        
        except KeyboardInterrupt:
            print("\n\nInterrupted by user. Goodbye!")
            break
        except Exception as e:
            print(f"\nError: {e}")
            traceback.print_exc()

if __name__ == "__main__":
    main()